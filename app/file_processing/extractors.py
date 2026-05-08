from __future__ import annotations

from html import unescape
from html.parser import HTMLParser
from pathlib import Path
import re

from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation


class TextExtractionError(Exception):
    pass


class HTMLVisibleTextParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._ignored_tag_depth = 0
        self.parts: list[str] = []

    def handle_starttag(self, tag: str, attrs) -> None:
        if tag.lower() in {"script", "style", "noscript"}:
            self._ignored_tag_depth += 1
        if tag.lower() in {"p", "br", "div", "section", "article", "li", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")

    def handle_endtag(self, tag: str) -> None:
        if tag.lower() in {"script", "style", "noscript"} and self._ignored_tag_depth:
            self._ignored_tag_depth -= 1
        if tag.lower() in {"p", "div", "section", "article", "li", "h1", "h2", "h3", "h4"}:
            self.parts.append("\n")

    def handle_data(self, data: str) -> None:
        if self._ignored_tag_depth:
            return
        if data.strip():
            self.parts.append(data)

    def text(self) -> str:
        return unescape(" ".join(self.parts))


def extract_text_from_txt(file_path: Path) -> str:
    raw_text = file_path.read_text(encoding="utf-8", errors="ignore")
    if looks_like_html(raw_text):
        return extract_visible_text_from_html(raw_text)
    return raw_text


def looks_like_html(text: str) -> bool:
    sample = text[:5000].lower()
    return bool(re.search(r"<!doctype\s+html|<html[\s>]|<body[\s>]|</(div|p|span|script|style|html)>", sample))


def extract_visible_text_from_html(html: str) -> str:
    parser = HTMLVisibleTextParser()
    parser.feed(html)
    return parser.text()


def extract_text_from_pdf(file_path: Path) -> str:
    reader = PdfReader(str(file_path))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n".join(pages)


def extract_text_from_docx(file_path: Path) -> str:
    document = DocxDocument(str(file_path))
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    return "\n".join(lines)


def extract_text_from_pptx(file_path: Path) -> str:
    presentation = Presentation(str(file_path))
    lines: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                lines.append(shape.text)
    return "\n".join(lines)


def clean_text(text: str) -> str:
    normalized = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.strip() for line in normalized.split("\n")]
    compact = "\n".join(line for line in lines if line)
    return compact.strip()


def extract_text_from_file(file_path: Path) -> str:
    suffix = file_path.suffix.lower().replace(".", "")
    if suffix == "txt":
        text = extract_text_from_txt(file_path)
    elif suffix == "pdf":
        text = extract_text_from_pdf(file_path)
    elif suffix == "docx":
        text = extract_text_from_docx(file_path)
    elif suffix == "pptx":
        text = extract_text_from_pptx(file_path)
    else:
        raise TextExtractionError(f"Unsupported extension: {suffix}")

    text = clean_text(text)
    if not text:
        raise TextExtractionError("Не удалось извлечь текст из файла")
    return text
